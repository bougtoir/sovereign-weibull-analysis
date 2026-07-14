#!/usr/bin/env python3
"""
Generate editable PPTX with 5 new revision figures (1 per slide).
Widescreen 13.333 x 7.5 inches.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(SCRIPT_DIR, '..', 'figures')
OUT_DIR = os.path.join(SCRIPT_DIR, '..', 'manuscript')

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FIGURES = [
    {
        'file': 'fig_model_selection.png',
        'title': 'Figure 9. Model Selection: AIC/BIC Comparison',
        'caption': (
            'Model selection comparison across Weibull, log-normal, gamma, and exponential '
            'distributions for all 31 polities. Left: ΔAIC relative to best-fitting model. '
            'Right: ΔBIC relative to best-fitting model. Filled circles = best model; '
            'open squares = competitive alternatives (Δ < 2). Vertical dashed lines at Δ = 2 and Δ = 10.'
        ),
    },
    {
        'file': 'fig_papacy_comparison.png',
        'title': 'Figure 10. Papacy: Pre-1000 CE vs Post-1000 CE',
        'caption': (
            '(A) Kaplan-Meier survival functions with fitted Weibull curves for pre-1000 CE (red, N=140) '
            'and post-1000 CE (blue, N=121). (B) Weibull probability plots (linearised ln-ln). '
            '(C) Reign duration histograms with fitted Weibull PDFs. '
            'Pre-1000 k=1.092 (near-exponential); post-1000 k=1.163 (mildly increasing hazard).'
        ),
    },
    {
        'file': 'fig_roman_saleh_comparison.png',
        'title': 'Figure 11. Roman Empire: Single vs Mixture Weibull',
        'caption': (
            '(A) Probability density functions: data histogram (grey), single Weibull (blue, k=0.880), '
            'two-component mixture (red), with individual components shown. '
            '(B) Hazard functions: single model shows monotonically decreasing hazard; '
            'mixture produces bathtub curve consistent with Saleh (2019).'
        ),
    },
    {
        'file': 'fig_s1_per_polity_fits.png',
        'title': 'Supplementary Figure S1. Per-Polity Empirical vs Fitted Distributions',
        'caption': (
            'Empirical Kaplan-Meier survival curves (black step functions) overlaid with fitted '
            'Weibull (blue solid), log-normal (red dashed), and gamma (green dotted) distributions '
            'for all 31 polities. Each panel shows polity name, sample size N, best-fitting model '
            'by AIC, and Weibull ΔAIC.'
        ),
    },
    {
        'file': 'fig_sample_size_analysis.png',
        'title': 'Figure 12. Sample Size Assessment',
        'caption': (
            '(A) 95% CI width vs sample size N for all 31 polities (red = CI crosses k=1; '
            'blue = does not). Vertical dashed line at N=15 indicates minimum threshold for '
            'qualitative interpretation. (B) Forest plot of 95% CIs ordered by sample size, '
            'with k=1 reference line.'
        ),
    },
]


def generate_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for fig_info in FIGURES:
        slide = prs.slides.add_slide(blank_layout)
        fig_path = os.path.join(FIG_DIR, fig_info['file'])

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = fig_info['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Image
        if os.path.exists(fig_path):
            # Calculate proportional size
            max_w = Inches(12.0)
            max_h = Inches(5.2)

            from PIL import Image
            with Image.open(fig_path) as img:
                img_w, img_h = img.size
            aspect = img_w / img_h
            if aspect > (max_w / max_h):
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)

            left = (SLIDE_W - w) // 2
            top = Inches(0.9)
            slide.shapes.add_picture(fig_path, left, top, w, h)
        else:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            tf2 = txBox2.text_frame
            tf2.paragraphs[0].text = f'[Image not found: {fig_info["file"]}]'
            tf2.paragraphs[0].font.size = Pt(14)

        # Caption
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.3), Inches(1.0))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = fig_info['caption']
        p3.font.size = Pt(11)
        p3.font.italic = True
        p3.alignment = PP_ALIGN.LEFT

    out_path = os.path.join(OUT_DIR, 'hassc_revision_figures.pptx')
    prs.save(out_path)
    print(f"Saved PPTX to: {out_path}")
    return out_path


if __name__ == '__main__':
    generate_pptx()
